"""Build the NoC performance-probe deck on top of the company master template.

Opens 20260608週報.pptx (corporate master / layouts / logo), drops its content
slides, and adds the performance-probe slides built from native shapes and tables
so everything stays editable and inherits the corporate look.

Run:  py -3 docs/build_perf_probe_slides.py
Out:  docs/performance-probe.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "20260608週報.pptx")
OUT = os.path.join(HERE, "performance-probe.pptx")

# Corporate theme palette (from theme4.xml)
DARK   = RGBColor(0x44, 0x54, 0x6A)
BLUE   = RGBColor(0x5B, 0x9B, 0xD5)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN  = RGBColor(0x70, 0xAD, 0x47)
GRAY   = RGBColor(0xA5, 0xA5, 0xA5)
LIGHT  = RGBColor(0xE7, 0xE6, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = 13.333, 7.5  # slide size in inches

prs = Presentation(TEMPLATE)

# ---- drop existing content slides, keep masters/layouts -------------------
# drop_rel removes the relationship so the orphan slide parts are not
# serialized (otherwise new slides collide with the old slideN.xml names).
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)

L_TITLE = prs.slide_layouts[4]   # title + subtitle
L_ONLY  = prs.slide_layouts[8]   # title only (free canvas below)


def title_only(title):
    s = prs.slides.add_slide(L_ONLY)
    s.shapes.title.text = title
    return s


def box(s, x, y, w, h, text="", fill=BLUE, line=None, tcol=WHITE,
        size=13, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        align=PP_ALIGN.CENTER):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold if i == 0 else False
        r.font.color.rgb = tcol
        if i == 0 and len(lines) > 1:
            r.font.bold = True
    return sp


def textbox(s, x, y, w, h, runs, size=13, tcol=DARK, align=PP_ALIGN.LEFT,
            bullet=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(4)
        r = p.add_run(); r.text = (("• " + ln) if bullet else ln)
        r.font.size = Pt(size); r.font.color.rgb = tcol
    return tb


def arrow(s, x, y, w, h, fill=GRAY):
    sp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def simple_table(s, x, y, w, rowh, data, colw, header_fill=DARK, size=12):
    rows, cols = len(data), len(data[0])
    h = rowh * rows
    gt = s.shapes.add_table(rows, cols, Inches(x), Inches(y),
                            Inches(w), Inches(h)).table
    for ci, cw in enumerate(colw):
        gt.columns[ci].width = Inches(cw)
    for ri in range(rows):
        gt.rows[ri].height = Inches(rowh)
        for ci in range(cols):
            c = gt.cell(ri, ci)
            c.margin_left = Inches(0.08); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(data[ri][ci])
            r.font.size = Pt(size)
            if ri == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
                r.font.bold = True; r.font.color.rgb = WHITE
            else:
                c.fill.solid()
                c.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = DARK
    return gt


# ============================ Slide 1: Title ===============================
s = prs.slides.add_slide(L_TITLE)
s.shapes.title.text = "NoC Performance Probe"
for ph in s.placeholders:
    if ph.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
        ph.text = "Two-view performance monitoring on the running design"

# ============================ Slide 2: Overview ============================
s = title_only("Overview")
box(s, 0.7, 1.7, 5.7, 2.2,
    "Profile view\nAggregate counters and per-class latency for each "
    "interface.\nMinimum = best observed latency.\nSpread to maximum = "
    "contention from other traffic.", fill=BLUE, size=14)
box(s, 6.9, 1.7, 5.7, 2.2,
    "Trace view\nOne record per completed transaction.\nUsed to drill into "
    "the transactions behind a profile outlier.", fill=ORANGE, size=14)
textbox(s, 0.7, 4.3, 11.9, 2.4, [
    "Non-intrusive: the monitors only observe. The design runs the same "
    "cycles with them on or off.",
    "One run produces one perf.json and one on-screen summary.",
    "The best-case reference is the smallest latency seen in that same run, "
    "not a separate idle pass.",
], size=15, bullet=True)

# ============================ Slide 3: Placement ===========================
s = title_only("Probe placement")
chain = [("AXI\nmaster", DARK), ("NMU", BLUE), ("Router", BLUE),
         ("Router", BLUE), ("NSU", BLUE), ("AXI\nmemory", DARK)]
x, y, bw, bh, gap = 0.55, 3.0, 1.6, 1.0, 0.42
xs = []
for i, (label, col) in enumerate(chain):
    bx = x + i * (bw + gap)
    xs.append(bx)
    box(s, bx, y, bw, bh, label, fill=col, size=13)
    if i < len(chain) - 1:
        arrow(s, bx + bw + 0.06, y + bh / 2 - 0.12, gap - 0.12, 0.24)
# link monitor tag between the two routers (index 2,3)
linkx = xs[2] + bw + 0.02
box(s, linkx - 0.15, y - 0.95, gap + 0.3, 0.55, "link\nmonitor",
    fill=GREEN, size=10)
arrow(s, linkx, y + bh / 2 - 0.12, gap - 0.12, 0.24, fill=GREEN)
# AXI monitor tags on manager edge (master->NMU) and memory edge (NSU->memory)
box(s, xs[0] + bw - 0.1, y + bh + 0.25, 1.7, 0.55, "AXI monitor\n(manager edge)",
    fill=ORANGE, size=10)
box(s, xs[4] + bw - 0.1, y + bh + 0.25, 1.7, 0.55, "AXI monitor\n(memory edge)",
    fill=ORANGE, size=10)
# router queue sampling tag
box(s, xs[2], y + bh + 0.25, 2.0, 0.55, "queue sampling\n(per cycle)",
    fill=GRAY, size=10)
textbox(s, 0.55, 1.55, 12.0, 0.9, [
    "Four tap types, all input-only: an AXI monitor on each manager and "
    "memory edge, a monitor on each link, and per-cycle router-queue sampling.",
    "Shown as one datapath; the 2-node testbench mirrors it.",
], size=13, bullet=True)

# ============================ Slide 4: Operation ===========================
s = title_only("Operation")
ow, oh, oy = 3.7, 2.3, 1.9
specs = [("INPUT", "Tap interface, link, and\nrouter-queue wires.\nEvery tap is "
          "read-only.", BLUE),
         ("COMPUTE", "On each clock, detect the\nhandshakes and update the\n"
          "counters (one set per monitor).", ORANGE),
         ("OUTPUT", "Derive throughput and\nper-class latency.\nWrite perf.json "
          "and the summary.", GREEN)]
ox = 0.55
oxs = []
for i, (h, body, col) in enumerate(specs):
    bx = ox + i * (ow + 0.55)
    oxs.append(bx)
    box(s, bx, oy, ow, oh, h + "\n" + body, fill=col, size=14)
    if i < 2:
        arrow(s, bx + ow + 0.08, oy + oh / 2 - 0.18, 0.40, 0.36)
textbox(s, 0.55, 4.7, 12.2, 2.0, [
    "Latency: AWVALID & AWREADY  →  BVALID & BREADY (write); "
    "ARVALID & ARREADY  →  RVALID & RREADY & RLAST (read).",
    "Idle: data VALID held without READY.   Link stall: credit == 0.",
    "Each start cycle is held in a per-id queue so writes and reads stay "
    "matched while several are in flight.",
], size=14, bullet=True)

# ============================ Slide 5: Profile metrics =====================
s = title_only("Profile view: event counting")
data = [
    ["Metric", "Definition"],
    ["Write / Read Transaction Count", "Completed write or read transactions."],
    ["Write / Read Byte Count", "Bytes written or read (basis for throughput)."],
    ["Write Latency", "AWVALID & AWREADY  →  BVALID & BREADY."],
    ["Read Latency", "ARVALID & ARREADY  →  RVALID & RREADY & RLAST."],
    ["Slave Write Idle Cycle Count", "Clocks WVALID is held without WREADY."],
    ["Master Read Idle Cycle Count", "Clocks RVALID is held without RREADY."],
    ["Outstanding", "Peak accepted-but-not-completed transactions."],
    ["Link / Router", "Flit count, stall (credit == 0), queue occupancy peak."],
]
simple_table(s, 0.55, 1.6, 12.2, 0.52, data, [4.3, 7.9], size=12)
textbox(s, 0.55, 6.5, 12.2, 0.7, [
    "Derived: throughput, per-class min / mean / max, and a fixed-range "
    "latency distribution.",
], size=13, bullet=True)

# ============================ Slide 6: Trace + composition =================
s = title_only("Trace view and latency composition")
textbox(s, 0.55, 1.55, 12.0, 0.9, [
    "Trace: one JSON record per completed transaction "
    "(direction, source, destination, accept and complete cycle, latency, "
    "bytes). Not printed in the summary.",
], size=13, bullet=True)
# stacked bars: 1 cycle = SCALE inches
SCALE = 0.36
bx0, lblw = 2.0, 1.5
segs_w = [("NI", 10, BLUE), ("Router", 12, ORANGE), ("Slave", 3, GREEN),
          ("Shell", 2, GRAY)]
segs_r = [("NI", 10, BLUE), ("Router", 12, ORANGE), ("Slave", 2, GREEN),
          ("Shell", 4, GRAY)]
for bi, (name, segs, total) in enumerate([("Write = 27 cyc", segs_w, 27),
                                          ("Read = 28 cyc", segs_r, 28)]):
    by = 3.1 + bi * 1.4
    textbox(s, 0.55, by + 0.12, lblw, 0.6, [name], size=13)
    cx = bx0
    for lab, cyc, col in segs:
        w = cyc * SCALE
        box(s, cx, by, w, 0.75, f"{lab}\n{cyc}", fill=col, size=11,
            shape=MSO_SHAPE.RECTANGLE)
        cx += w
textbox(s, 0.55, 6.2, 12.2, 1.0, [
    "Measured directly: round-trip total and slave service. NI and router "
    "come from their pipeline depth; shell is the remainder. One measured "
    "instance; regenerate before reuse.",
], size=12, bullet=True)

# ---- geometry sanity check ------------------------------------------------
def emu_in(v):
    return v / 914400.0
warn = 0
for idx, sl in enumerate(prs.slides, 1):
    for sp in sl.shapes:
        try:
            l, t, w, h = emu_in(sp.left), emu_in(sp.top), emu_in(sp.width), emu_in(sp.height)
        except Exception:
            continue
        if l < -0.05 or t < -0.05 or l + w > SW + 0.05 or t + h > SH + 0.05:
            warn += 1
            print(f"  [warn] slide{idx} shape out of bounds: "
                  f"x={l:.2f} y={t:.2f} w={w:.2f} h={h:.2f}")
print(f"slides: {len(prs.slides._sldIdLst)}  out-of-bounds warnings: {warn}")
prs.save(OUT)
print("saved:", OUT)
